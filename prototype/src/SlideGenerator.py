# primary author: Pramod

from pptx import Presentation
from pptx.util import Inches, Pt

def create_title_slide(prs, inp_title = "", inp_subtitle = ""):
	title_slide_layout = prs.slide_layouts[0]
	slide = prs.slides.add_slide(title_slide_layout)
	title = slide.shapes.title
	subtitle = slide.placeholders[1]

	title.text = inp_title
	subtitle.text = inp_subtitle

def setFooter(prs, slide_no,footer_text):
        slide = prs.slides[slide_no]
        left = Inches(2.8)
        width = Inches(3.7)
        top = Inches(6.7)
        height = Inches(1)
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_box_frame = text_box.text_frame
        footer_paragraph = text_box_frame.add_paragraph()
        footer_paragraph.text = footer_text
        footer_paragraph.font.bold = True
        footer_paragraph.font.size = Pt(13)

def add_bullet_slide(prs, title, array_of_bullet_sentences):
	bullet_slide_layout = prs.slide_layouts[1]
	slide = prs.slides.add_slide(bullet_slide_layout)
	shapes = slide.shapes
	title_shape = shapes.title
	body_shape = shapes.placeholders[1]
	title_shape.text = title
	tf = body_shape.text_frame
	if(len(array_of_bullet_sentences) > 0):
		tf.text = array_of_bullet_sentences[0]
		for bullet in array_of_bullet_sentences[1:]:
			p = tf.add_paragraph()
			p.text = bullet
			p.font.size = Pt(18)

def add_text_slide(prs,text_array,title = ""):
	blank_slide_layout = prs.slide_layouts[6]
	slide = prs.slides.add_slide(blank_slide_layout)

	# left = top = width = height = Inches(1)
	left = Inches(0.5)
	width = Inches(5)
	top = Inches(1)
	height = Inches(5)
	txBox = slide.shapes.add_textbox(left, top, width, height)
	tf = txBox.text_frame
	p = tf.add_paragraph()
	p.text = title
	p.font.size = Pt(35)

	p = tf.add_paragraph()
	p.text = ""
	p.font.size = Pt(20)

	for sentence in text_array:
		p = tf.add_paragraph()
		p.text = sentence
		p.font.size = Pt(20)

def add_image(prs, slide_no, image_path):
	slide = prs[slide_no]
	left = top = Inches(1)
	pic = slide.shapes.add_picture(image_path, left, top)


def setLogo(prs,slide_no,logo_path):
    slide = prs.slides[slide_no]
    left = Inches(0.2)
    top = Inches(0.2)
    slide.shapes.add_picture(logo_path,left,top)

def create_presentation(file_name):
	prs = Presentation()
	create_title_slide(prs, 'Hello, World!', 'python-pptx was here!')
	# setFooter(prs, 2, 'PESIT SUCKS')
	# setLogo(prs,2,'logo.png')
	bullet_title = 'BULLETS'
	bullets = ['bullet one', 'bullet two']
	add_bullet_slide(prs, bullet_title, bullets)
	add_text_slide(prs, ['jksdhflkadhsofhsakdhbf','asdfsfd'], 'TEXT HERE')
	prs.save(file_name+ '.pptx')

# create_presentation('test')