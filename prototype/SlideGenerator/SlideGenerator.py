from pptx import Presentation
from pptx.util import Inches, Pt
from nltk import word_tokenize
import re
from textblob import TextBlob
from gensim.summarization import keywords
from nltk.stem import WordNetLemmatizer

class SlideGenerator:

    def __init__(self):
        self.prs = Presentation()
    
    def lemmatize_word(self, word, pos='n'):
        wordnet_lemmatizer = WordNetLemmatizer()
        return wordnet_lemmatizer.lemmatize(word, pos=pos)

    def get_keywords(self, text, ratio=0.1):
        return keywords(text, ratio=ratio, split=True)

    def get_nouns(self, text):
        blob = TextBlob(text)
        return list(blob.noun_phrases)
    
    def get_bullet_title(self, sentences):
        text = ' '.join(sentences)
        keywords = self.get_keywords(text, ratio=0.1)
        nouns = self.get_nouns(text)
        if len(nouns) == 0 or len(keywords) == 0:
            return "<Please Fill in an appropriate title>"
        likely_titles = list(set(keywords) and set(nouns))
        if len(likely_titles) == 1:
            return self.lemmatize_word(likely_titles[0])
        elif len(likely_titles) > 1:
            return self.lemmatize_word(likely_titles[0])

    def get_cleaned_bullets(self, sentences):
        c_sentences = ['']
        for sentence in sentences:
            words = word_tokenize(sentence)
            if len(words) in xrange(2, 30):
                sentence = re.sub(r'^\W+', '', sentence)
                sentence = re.sub(r'^[0-9\.,?:;!\)\}\]]*', '', sentence)
                sentence = re.sub(r'[0-9]*$', '', sentence)
                sentence = re.sub(r'\W+$', '', sentence)
                sentence = re.sub(r'\W+$', '', sentence)

                sentence = re.sub(r'- ', '', sentence)

                sentence = re.sub(r'[\(\{\[][\w+ \t\n,:;?]*$', '', sentence)

                sentence = sentence.strip()
                if len(sentence) > 1:
                    sentence = sentence[0].upper()+sentence[1:]
                c_sentences.append(sentence)
            else:
                #ignore the sentence
                pass
        return c_sentences

    def create_title_slide(self, inp_title="", inp_subtitle=""):
        title_slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = inp_title
        subtitle.text = inp_subtitle

    def set_footer(self, slide_no, footer_text):
            slide = self.prs.slides[slide_no]
            left = Inches(2.8)
            width = Inches(3.7)
            top = Inches(6.7)
            height = Inches(1)
            text_box = slide.shapes.add_textbox(left, top, width, height)
            text_box_frame = text_box.text_frame
            footer_paragraph = text_box_frame.add_paragraph()
            footer_paragraph.text = footer_text
            footer_paragraph.font.bold = True
            footer_paragraph.font.size = Pt(13)

    def add_bullet_slide(self, title, array_of_bullet_sentences):
        bullet_slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = title
        tf = body_shape.text_frame
        if len(array_of_bullet_sentences) > 0:
            tf.text = array_of_bullet_sentences[0]
            for bullet in array_of_bullet_sentences[1:]:
                p = tf.add_paragraph()
                p.text = bullet
                p.font.size = Pt(18)

    def add_text_slide(self, text_array,title = ""):
        blank_slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_slide_layout)

        # left = top = width = height = Inches(1)
        left = Inches(0.5)
        width = Inches(5)
        top = Inches(1)
        height = Inches(5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = title
        p.font.size = Pt(35)

        p = tf.add_paragraph()
        p.text = ""
        p.font.size = Pt(20)

        for sentence in text_array:
            p = tf.add_paragraph()
            p.text = sentence
            p.font.size = Pt(20)

    def add_image(self, slide_no, image_path):
        slide = self.prs[slide_no]
        left = top = Inches(1)
        pic = slide.shapes.add_picture(image_path, left, top)

    def set_logo(self, slide_no, logo_path):
        if logo_path != '':
            slide = self.prs.slides[slide_no]
            left = Inches(0.2)
            top = Inches(0.2)
            slide.shapes.add_picture(logo_path, left, top)

    def create_presentation(self, output_file, title='', sub_title='', footer='', logo='logo.png', contents=[]):
        self.create_title_slide(title, sub_title)
        for i in range(0, len(contents), 5):
            bullets = contents[i:i+5]
            bullets.insert(0, "")
            bullet_title = self.get_bullet_title(bullets).title()
            bullets = self.get_cleaned_bullets(bullets)
            self.add_bullet_slide(bullet_title, bullets)
            self.set_logo((i/5)+1, logo)
            self.set_footer((i/5)+1, footer)
        self.prs.save(output_file + '.pptx')