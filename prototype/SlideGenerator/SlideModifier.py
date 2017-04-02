from pptx import Presentation
from pptx.util import Inches, Pt

"""
there are already APIs for levelling so not a problem
will look into these :
http://pbpython.com/creating-powerpoint.html
http://python-pptx.readthedocs.io/en/latest/user/quickstart.html
"""

def setFooter(from_slide_number,to_slide_number,slide_ref,footer_text):
    for slide_number in range(from_slide_number,to_slide_number+1):
        # position reference from top left corner
        slide = slide_ref[slide_number]
        left = Inches(2.8)
        width = Inches(3.7)
        top = Inches(6.7)
        height = Inches(1)
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_box_frame = text_box.text_frame
        # need to look into removal of unnecessarily new line character created in text box
        footer_paragraph = text_box_frame.add_paragraph()
        footer_paragraph.text = footer_text
        footer_paragraph.font.bold = True
        footer_paragraph.font.size = Pt(10)
    return slide_ref

def setLogo(from_slide_number,to_slide_number,slide_ref,logo_path):
    for slide_number in range(from_slide_number,to_slide_number+1):
        slide = slide_ref[slide_number]
        left = Inches(8.6)
        top = Inches(0.2)
        slide.shapes.add_picture(logo_path,left,top)
    return slide_ref


presentation_object =  Presentation("procs.pptx")
setFooter(0,2,presentation_object.slides,"Department of ISE,PESIT")
setLogo(0,5,presentation_object.slides,"logo.png")
presentation_object.save("procs-mod.pptx")
